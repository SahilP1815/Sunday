import os
import re
import sys
import json
import subprocess
import urllib.parse
import tempfile
import requests
from io import BytesIO
from pathlib import Path
from PIL import Image
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def safe_print(msg: str):
    """Print message safely, handling potential encoding errors on Windows console."""
    try:
        print(msg)
    except UnicodeEncodeError:
        try:
            encoding = sys.stdout.encoding or 'ascii'
            print(msg.encode(encoding, errors='replace').decode(encoding))
        except Exception:
            try:
                print(msg.encode('ascii', errors='replace').decode('ascii'))
            except Exception:
                pass

# Define a premium, modern design theme palette
COLOR_BG = RGBColor(15, 23, 42)         # Slate 900
COLOR_CARD = RGBColor(30, 41, 59)       # Slate 800
COLOR_TITLE = RGBColor(56, 189, 248)    # Sky Blue
COLOR_BODY = RGBColor(226, 232, 240)     # Slate 200
COLOR_DESC = RGBColor(203, 213, 225)     # Slate 300
COLOR_SUBTITLE = RGBColor(148, 163, 184) # Slate 400

SARVAM_CHAT_URL    = "https://api.sarvam.ai/v1/chat/completions"
SARVAM_PPT_MODEL   = "sarvam-m"
SARVAM_PPT_FALLBACK = "sarvam-30b"

class PPTGeneratorSkill:
    def __init__(self, gemini_api_key: str = None, sarvam_api_key: str = None):
        self.gemini_api_key = gemini_api_key
        self.sarvam_api_key = sarvam_api_key
        if gemini_api_key:
            genai.configure(api_key=gemini_api_key)

    def parse_command(self, text: str) -> dict:
        """Detect intent and extract PPT topic and slide count."""
        text_lower = text.lower().strip()
        
        # Regex to detect if command is asking to create a presentation
        match_intent = re.search(r"\b(create|make|generate|build|design|write)\s+(?:a\s+)?(?:new\s+)?(ppt|powerpoint|presentation|slides)\b", text_lower)
        if not match_intent and "ppt" not in text_lower and "presentation" not in text_lower:
            return {"intent": None, "topic": None, "slide_count": 0}
            
        # Extract number of slides if mentioned
        slide_match = re.search(r"(\d+)\s*(slide|page|ppt)", text_lower)
        slide_count = int(slide_match.group(1)) if slide_match else 5  # default to 5 slides
        
        # Extract the topic by removing typical command verbs and stop words
        clean_text = re.sub(
            r"\b(create|make|generate|build|design|write|a|an|new|ppt|powerpoint|presentation|slides|on|about|of|for|topic)\b",
            "",
            text_lower
        )
        # Strip punctuation and double spaces
        clean_text = re.sub(r"\s+", " ", clean_text).strip()
        topic = clean_text.title() if clean_text else "General Knowledge"
        
        return {
            "intent": "create_ppt",
            "topic": topic,
            "slide_count": slide_count
        }

    def _get_wikipedia_context(self, topic: str) -> str:
        search_url = "https://en.wikipedia.org/w/api.php"
        search_params = {
            "action": "query",
            "list": "search",
            "srsearch": topic,
            "utf8": "",
            "format": "json",
            "srlimit": 1
        }
        headers = {
            "User-Agent": "SundayAssistant/1.0 (contact: admin@example.com)"
        }
        try:
            safe_print(f"[PPT Skill] Searching Wikipedia for topic: '{topic}'...")
            res = requests.get(search_url, params=search_params, headers=headers, timeout=5)
            if res.status_code == 200:
                results = res.json().get("query", {}).get("search", [])
                if results:
                    page_title = results[0].get("title")
                    safe_print(f"[PPT Skill] Fetching summary for Wikipedia page: '{page_title}'...")
                    summary_params = {
                        "action": "query",
                        "prop": "extracts",
                        "exintro": "",
                        "explaintext": "",
                        "redirects": 1,
                        "format": "json",
                        "titles": page_title
                    }
                    res_sum = requests.get(search_url, params=summary_params, headers=headers, timeout=5)
                    if res_sum.status_code == 200:
                        pages = res_sum.json().get("query", {}).get("pages", {})
                        for page_id, page in pages.items():
                            extract = page.get("extract", "")
                            if extract:
                                return f"Wikipedia Page Title: {page_title}\nSummary Extract: {extract}"
        except Exception as e:
            safe_print(f"[PPT Skill] Error fetching Wikipedia context: {e}")
        return ""

    def _get_web_context(self, topic: str) -> str:
        """Search the web via the Express backend server for context on the topic."""
        try:
            safe_print(f"[PPT Skill] Searching the web for topic: '{topic}'...")
            
            headers = {}
            config_file = Path(__file__).parent / "config.json"
            if config_file.exists():
                try:
                    config_data = json.loads(config_file.read_text())
                    tavily_key = config_data.get("tavily_api_key", "").strip()
                    if tavily_key:
                        headers["x-tavily-api-key"] = tavily_key
                except Exception:
                    pass

            response = requests.get(
                "http://localhost:3001/api/search",
                params={"q": topic},
                headers=headers,
                timeout=12
            )
            if response.status_code == 200:
                data = response.json()
                results = data.get("results", [])
                if not results:
                    return ""
                
                formatted = []
                for r in results[:5]:
                    title = r.get("title", "No Title")
                    snippet = r.get("snippet", "")
                    url = r.get("url", "")
                    formatted.append(f"Title: {title}\nSnippet: {snippet}\nURL: {url}\n")
                
                context = "\n".join(formatted)
                safe_print(f"[PPT Skill] Successfully retrieved web search context. {len(results)} results.")
                return context
            else:
                safe_print(f"[PPT Skill] Search API returned status: {response.status_code}")
        except Exception as e:
            safe_print(f"[PPT Skill] Error fetching web search context: {e}")
        return ""

    def _build_prompt(self, topic: str, slide_count: int, reference_context: str) -> str:
        """Build the shared content-generation prompt used by both Sarvam and Gemini."""
        prompt = f"""\
You are a professional PowerPoint content designer. Create a highly detailed and comprehensive structured outline for a presentation with {slide_count} slides on the topic/prompt: "{topic}".
If the prompt specifies particular focus areas, sub-topics, or details to include, you MUST design the presentation around those specifications. If no specific details are mentioned, generate a standard, well-rounded introduction and structured overview of the topic.

For the overall presentation, extract a short, clear keyword-based title (e.g., "Machine Learning Basics" instead of "a presentation on the history of machine learning and its future applications in modern healthcare").

CRITICAL IMAGE QUERY REQUIREMENT:
All image queries generated MUST represent high-quality, professional photography, clean digital designs, or clear conceptual diagrams.
DO NOT use obscure, niche, or overly specific proper nouns that Bing Image Search would fail to match (e.g. do not search for album names, local news events, or private entities directly). Instead, translate niche topics into generic, visually appealing concepts.
Examples:
- If the topic is about "MGK's Mainstream Sellout Album", search for "electric guitar rock concert stage neon lights pink background" or "vinyl record spinning".
- If the topic is about a niche programming library, search for "clean software code on computer screen" or "digital technology network nodes".
- Avoid search terms that might retrieve news screenshots, text-heavy pages, or memes.
"""
        if reference_context:
            prompt += f"\n\nHere is real-time web search and verified factual context about the topic:\n{reference_context}\nUse this information to ensure all slides, descriptions, and bullet points are factual, up-to-date, accurate, and contain relevant, detailed data."

        prompt += """

For the presentation, you MUST provide:
1. A "refined_topic": a short keyword-based title.
2. A "cover_image_query": a specific, highly relevant search term for a beautiful cover photo (e.g., "pop punk guitar concert stage lights" for an album topic).
3. A list of slides. For each slide, provide:
   - A compelling, clear "title".
   - A "description": a detailed paragraph of 2-3 sentences explaining the slide's theme in depth.
   - A list of 3-4 "bullets" that provide deep, specific, and structured sub-points. Each bullet point should be a complete, well-formed sentence.
   - An "image_query": a specific, descriptive search term for a stock photo representing the slide's theme. If no image is needed, leave it as an empty string "".

Return ONLY a JSON object representing the presentation.
Do not include markdown tags like ```json or ```. Return pure JSON.

JSON Structure:
{
  "refined_topic": "Short Keyword Title",
  "cover_image_query": "high-quality photo search query for cover slide",
  "slides": [
    {
      "title": "Slide Title",
      "description": "Detailed explanatory paragraph for the slide.",
      "bullets": [
        "Detailed bullet point 1 with explanation.",
        "Detailed bullet point 2 with explanation.",
        "Detailed bullet point 3 with explanation."
      ],
      "image_query": "specific search query for a relevant high-quality photo or empty string"
    }
  ]
}"""
        return prompt

    def _parse_json_response(self, raw: str, topic: str) -> dict:
        """Strip markdown fences and parse the JSON response."""
        clean = raw.replace("```json", "").replace("```", "").strip()
        data = json.loads(clean)
        if isinstance(data, list):
            return {
                "refined_topic": topic,
                "cover_image_query": f"{topic} concept",
                "slides": data
            }
        return data

    def _generate_with_sarvam(self, prompt: str, topic: str) -> dict:
        """Call Sarvam AI chat completions endpoint. Raises on failure."""
        headers = {
            "api-subscription-key": self.sarvam_api_key,
            "Content-Type": "application/json"
        }
        payload = {
            "model": SARVAM_PPT_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7
        }
        response = requests.post(SARVAM_CHAT_URL, json=payload, headers=headers, timeout=30)
        if response.status_code == 404:
            # Model variant not found — retry with base model
            payload["model"] = SARVAM_PPT_FALLBACK
            response = requests.post(SARVAM_CHAT_URL, json=payload, headers=headers, timeout=30)
        if response.status_code not in (200, 201):
            raise Exception(f"Sarvam API {response.status_code}: {response.text[:200]}")
        raw = response.json()["choices"][0]["message"]["content"]
        return self._parse_json_response(raw, topic)

    def _generate_with_gemini(self, prompt: str, topic: str) -> dict:
        """Call Gemini API. Raises on failure."""
        model = genai.GenerativeModel("gemini-2.5-flash")
        response = model.generate_content(prompt)
        return self._parse_json_response(response.text, topic)

    def generate_content(self, topic: str, slide_count: int) -> dict:
        """Fetch structured slide content — Sarvam primary, Gemini fallback, templates last resort."""
        wiki_context = self._get_wikipedia_context(topic)
        web_context = self._get_web_context(topic)
        
        combined_context = ""
        if wiki_context:
            combined_context += f"--- Wikipedia Context ---\n{wiki_context}\n\n"
        if web_context:
            combined_context += f"--- Web Search Context ---\n{web_context}\n"
            
        if not combined_context:
            combined_context = "No direct web search reference found. Use your own internal knowledge base."
            
        prompt = self._build_prompt(topic, slide_count, combined_context)

        # 1. Try Sarvam AI (primary)
        if self.sarvam_api_key:
            try:
                print("[PPT Skill] Generating content with Sarvam AI...")
                data = self._generate_with_sarvam(prompt, topic)
                print("[PPT Skill] Sarvam AI content generation succeeded.")
                return data
            except Exception as e:
                print(f"[PPT Skill] Sarvam generation failed: {e}. Trying Gemini...")

        # 2. Try Gemini (fallback)
        if self.gemini_api_key:
            try:
                print("[PPT Skill] Generating content with Gemini...")
                data = self._generate_with_gemini(prompt, topic)
                print("[PPT Skill] Gemini content generation succeeded.")
                return data
            except Exception as e:
                print(f"[PPT Skill] Gemini generation failed: {e}. Using built-in templates.")

        # 3. Built-in template fallback
        print("[PPT Skill] No AI key available. Using fallback template content.")
        return self._get_fallback_content(topic, slide_count)

    def _get_fallback_content(self, topic: str, slide_count: int) -> dict:
        slides = []
        words = [w for w in topic.split() if w.lower() not in ["a", "an", "the", "presentation", "on", "about", "ppt", "slides"]]
        refined_title = " ".join(words[:4]).title() if words else "Overview"
        
        for i in range(1, slide_count + 1):
            slides.append({
                "title": f"Understanding {refined_title} - Part {i}",
                "description": f"This slide introduces the fundamental aspects of {refined_title}, outlining why it is significant in today's landscape. We examine core components, key challenges, and general methodologies used in the industry.",
                "bullets": [
                    f"Core foundational concept number {i} regarding {refined_title}.",
                    f"Practical applications and modern-day use cases within the sector.",
                    f"Important challenges, limitations, and the future outlook."
                ],
                "image_query": f"{refined_title} technology"
            })
        return {
            "refined_topic": refined_title,
            "cover_image_query": f"{refined_title} concept",
            "slides": slides
        }

    def _get_image_urls(self, query: str) -> list:
        url = f"https://www.bing.com/images/async?q={urllib.parse.quote(query)}&first=0&count=15"
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/115.0.0.0 Safari/537.36',
            'Accept-Language': 'en-US,en;q=0.9'
        }
        try:
            res = requests.get(url, headers=headers, timeout=10)
            if res.status_code == 200:
                urls = re.findall(r'murl&quot;:&quot;(http[^&]+?)&quot;', res.text)
                if not urls:
                    urls = re.findall(r'"murl":"(http[^"]+?)"', res.text)
                return list(dict.fromkeys(urls))
        except Exception as e:
            print(f"[PPT Skill] Error fetching image search for '{query}': {e}")
        return []

    def _crop_and_resize_image(self, img_data: bytes, target_w_in: float, target_h_in: float, dpi: int = 150) -> bytes:
        try:
            img = Image.open(BytesIO(img_data))
            target_w_px = int(target_w_in * dpi)
            target_h_px = int(target_h_in * dpi)
            
            target_ratio = target_w_px / target_h_px
            img_w, img_h = img.size
            img_ratio = img_w / img_h
            
            if img_ratio > target_ratio:
                new_w = int(img_h * target_ratio)
                left = (img_w - new_w) // 2
                img = img.crop((left, 0, left + new_w, img_h))
            else:
                new_h = int(img_w / target_ratio)
                top = (img_h - new_h) // 2
                img = img.crop((0, top, img_w, top + new_h))
                
            img = img.resize((target_w_px, target_h_px), Image.Resampling.LANCZOS)
            
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
                
            out_buf = BytesIO()
            img.save(out_buf, format="JPEG", quality=90)
            return out_buf.getvalue()
        except Exception as e:
            safe_print(f"[PPT Skill] Error cropping/resizing image: {e}")
            return None

    def _download_and_process_image(self, query: str, width_in: float, height_in: float, temp_dir: str, name_prefix: str) -> str:
        urls = self._get_image_urls(query)
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/115.0.0.0 Safari/537.36'
        }
        
        for url in urls[:5]:
            try:
                safe_print(f"[PPT Skill] Downloading for query '{query}': {url}")
                res = requests.get(url, headers=headers, timeout=5)
                if res.status_code == 200:
                    processed_data = self._crop_and_resize_image(res.content, width_in, height_in)
                    if processed_data:
                        save_path = os.path.join(temp_dir, f"{name_prefix}.jpg")
                        with open(save_path, "wb") as f:
                            f.write(processed_data)
                        return save_path
            except Exception as e:
                safe_print(f"[PPT Skill] Failed to process url {url}: {e}")
        return None

    def build_pptx_file(self, topic: str, slides_content) -> str:
        """Builds a styled PPTX file and returns the absolute file path."""
        # Normalize slides_content to dict format
        if isinstance(slides_content, list):
            refined_topic = topic
            slides_list = slides_content
        elif isinstance(slides_content, dict):
            refined_topic = slides_content.get("refined_topic", topic)
            slides_list = slides_content.get("slides", [])
        else:
            refined_topic = topic
            slides_list = []

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        temp_dir = tempfile.gettempdir()
        temp_image_paths = []
        
        # 1. Slide 1: Cover Slide (Blank layout)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_BG
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.0), Inches(4.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = refined_topic
        p.font.name = "Segoe UI"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE
        p.space_after = Pt(20)
        
        p_sub = tf.add_paragraph()
        p_sub.text = "Generated by Sunday AI Assistant\nPowered by python-pptx & Gemini"
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = COLOR_SUBTITLE
        
        # Download and add cover image
        cover_query = slides_content.get("cover_image_query", f"{refined_topic} concept") if isinstance(slides_content, dict) else f"{refined_topic} concept"
        cover_img_path = self._download_and_process_image(cover_query, 5.2, 5.9, temp_dir, "cover_slide")
        
        # Image card background
        frame = slide.shapes.add_shape(1, Inches(7.3), Inches(0.8), Inches(5.2), Inches(5.9))
        frame.fill.solid()
        frame.fill.fore_color.rgb = COLOR_CARD
        frame.line.fill.background()
        
        if cover_img_path and os.path.exists(cover_img_path):
            slide.shapes.add_picture(cover_img_path, Inches(7.4), Inches(0.9), Inches(5.0), Inches(5.7))
            temp_image_paths.append(cover_img_path)
        else:
            p_box = slide.shapes.add_textbox(Inches(7.4), Inches(0.9), Inches(5.0), Inches(5.7))
            p_tf = p_box.text_frame
            p_tf.word_wrap = True
            p_p = p_tf.paragraphs[0]
            p_p.text = refined_topic
            p_p.alignment = 1
            p_p.font.name = "Segoe UI"
            p_p.font.size = Pt(24)
            p_p.font.bold = True
            p_p.font.color.rgb = COLOR_SUBTITLE

        # 2. Add Content Slides
        for idx, slide_data in enumerate(slides_list):
            if slide_data.get("title").lower() == refined_topic.lower():
                continue
                
            slide = prs.slides.add_slide(blank_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG
            
            # Top thin accent line
            accent = slide.shapes.add_shape(1, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.05))
            accent.fill.solid()
            accent.fill.fore_color.rgb = COLOR_TITLE
            accent.line.fill.background()
            
            # Slide Title
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "Overview")
            p.font.name = "Segoe UI"
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLOR_TITLE
            
            # Check if slide has a valid image query
            image_query = slide_data.get("image_query", "").strip()
            img_path = None
            if image_query:
                img_name = f"content_slide_{idx}"
                img_path = self._download_and_process_image(image_query, 5.0, 4.6, temp_dir, img_name)
            
            # Determine if we should lay out in 2-column or full-width
            if img_path and os.path.exists(img_path):
                # 2-column layout (Left: Text, Right: Image)
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
                
                # Draw the card frame and image
                frame = slide.shapes.add_shape(1, Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.8))
                frame.fill.solid()
                frame.fill.fore_color.rgb = COLOR_CARD
                frame.line.fill.background()
                
                slide.shapes.add_picture(img_path, Inches(7.4), Inches(1.9), Inches(5.0), Inches(4.6))
                temp_image_paths.append(img_path)
            else:
                # Full-width text layout (since no image is requested/downloaded)
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
                
            # Style Text Content
            tf_c = content_box.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
            
            p_desc = tf_c.paragraphs[0]
            p_desc.text = slide_data.get("description", "No description provided.")
            p_desc.font.name = "Segoe UI"
            p_desc.font.size = Pt(15)
            p_desc.font.italic = True
            p_desc.font.color.rgb = COLOR_DESC
            p_desc.space_after = Pt(20)
            
            bullets = slide_data.get("bullets", [])
            for b_idx, bullet in enumerate(bullets):
                p_b = tf_c.add_paragraph()
                p_b.text = f"\u2022  {bullet}"
                p_b.font.name = "Segoe UI"
                p_b.font.size = Pt(14)
                p_b.font.color.rgb = COLOR_BODY
                p_b.space_after = Pt(12)

        # Save to user Documents folder
        docs_dir = os.path.join(os.path.expanduser("~"), "Documents")
        os.makedirs(docs_dir, exist_ok=True)
        
        safe_name = re.sub(r'[\\/*?:"<>|]', "", refined_topic).replace(" ", "_")
        file_path = os.path.join(docs_dir, f"{safe_name}_Presentation.pptx")
        
        prs.save(file_path)
        
        # Clean up temp images
        for path in temp_image_paths:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception:
                pass
                
        return file_path

        # Save to user Documents folder
        docs_dir = os.path.join(os.path.expanduser("~"), "Documents")
        os.makedirs(docs_dir, exist_ok=True)
        
        safe_name = re.sub(r'[\\/*?:"<>|]', "", topic).replace(" ", "_")
        file_path = os.path.join(docs_dir, f"{safe_name}_Presentation.pptx")
        
        prs.save(file_path)
        
        # Clean up temp images
        for path in temp_image_paths:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception:
                pass
                
        return file_path

    def launch_powerpoint(self, file_path: str):
        """Open the generated file in PowerPoint on Windows."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        if sys.platform == "win32":
            os.startfile(file_path)
        else:
            opener = "open" if sys.platform == "darwin" else "xdg-open"
            subprocess.Popen([opener, file_path])
